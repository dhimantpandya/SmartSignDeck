import pptx
prs = pptx.Presentation("Dhimant Pandya65.pptx")
print(f"Audit of Dhimant Pandya65.pptx ({len(prs.slides)} Slides):")

for i, slide in enumerate(prs.slides):
    possible_texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            possible_texts.append(shape.text.strip())
            
    content_summary = " | ".join(possible_texts)[:200]
    image_count = len([s for s in slide.shapes if s.shape_type == 13])
    print(f"Slide {i+1}: Images: {image_count} | Content: {content_summary}")
