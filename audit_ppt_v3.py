import pptx
prs = pptx.Presentation("Dhimant Pandya65.pptx")
print(f"Audit of Dhimant Pandya65.pptx ({len(prs.slides)} Slides):")

for i, slide in enumerate(prs.slides):
    # Try multiple ways to find a title or identifying text
    title = "Untitled"
    possible_texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            possible_texts.append(shape.text.strip())
            if shape == slide.shapes.title:
                title = shape.text.strip()
    
    # If no official title, use the first bit of text found
    if title == "Untitled" and possible_texts:
        title = possible_texts[0][:50]
        
    image_count = len([s for s in slide.shapes if s.shape_type == 13]) # 13 is Picture
    print(f"Slide {i+1}: {title} | Images: {image_count}")
