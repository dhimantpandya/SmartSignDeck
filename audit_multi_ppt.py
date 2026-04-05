import pptx
import sys

def audit(path):
    try:
        prs = pptx.Presentation(path)
        print(f"\nAudit of {path} ({len(prs.slides)} Slides):")
        for i, slide in enumerate(prs.slides):
            texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            summary = " | ".join(texts)[:150]
            print(f"Slide {i+1}: {summary}")
    except:
        print(f"Failed to audit {path}")

if __name__ == "__main__":
    audit("SmartSignDeck_Presentation_ULTIMATE_1775249086267.pptx")
    audit("SmartSignDeck_Presentation_ULTIMATE.pptx")
    audit("SmartSignDeck_Presentation_Final_Delivery.pptx")
