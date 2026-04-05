import pptx
try:
    prs = pptx.Presentation("Dhimant Pandya65.pptx")
    print("Audit of Dhimant Pandya65.pptx:")
    for i, slide in enumerate(prs.slides):
        title = "Untitled"
        if slide.shapes.title:
            title = slide.shapes.title.text
        print(f"Slide {i+1}: {title}")
except Exception as e:
    print(f"Error: {e}")
