import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = pptx.Presentation("Dhimant Pandya65.pptx")
print(f"Audit of Dhimant Pandya65.pptx ({len(prs.slides)} Slides):")

for i, slide in enumerate(prs.slides[:100]): # Limit to avoid excessive output but coverage is 98
    title = "Untitled"
    if slide.shapes.title:
        title = slide.shapes.title.text
    
    # Count images
    images = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    
    # Extract ALL text from the slide to identify its topic
    full_text = ""
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            full_text += shape.text.strip() + " | "
    
    print(f"Slide {i+1}: {title} | Images: {len(images)} | Text: {full_text[:150]}")
